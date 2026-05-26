#!/usr/bin/env python
"""Render v1 scene JSON to a PowerPoint file with python-pptx."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any


EMU_PER_INCH = 914400
PX_PER_INCH = 96


def px(value: float) -> int:
    return int(round(float(value) / PX_PER_INCH * EMU_PER_INCH))


def require_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore
        from pptx.util import Pt  # type: ignore
        from pptx.oxml import parse_xml  # type: ignore
        from pptx.oxml.ns import nsdecls, qn  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency
        raise SystemExit(
            "python-pptx is required for PPTX generation. Install it in the active Python "
            f"environment, then rerun this script. Import error: {exc}"
        )
    return Presentation, RGBColor, MSO_SHAPE, PP_ALIGN, MSO_ANCHOR, Pt, parse_xml, nsdecls, qn


def hex_to_rgb(color: str):
    color = color.strip().lstrip("#")
    if len(color) != 6:
        color = "000000"
    return tuple(int(color[index : index + 2], 16) for index in (0, 2, 4))


def resolve_path(scene_dir: Path, value: str) -> Path:
    path = Path(value)
    return path if path.is_absolute() else scene_dir / path


def add_text(slide: Any, element: dict[str, Any], helpers: tuple[Any, ...]) -> None:
    _, RGBColor, _, PP_ALIGN, MSO_ANCHOR, Pt, _, _, _ = helpers
    box = slide.shapes.add_textbox(px(element["x"]), px(element["y"]), px(element["width"]), px(element["height"]))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = str(element.get("content", ""))
    style = element.get("style", {})
    alignment = str(style.get("alignment", "left")).lower()
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(alignment, PP_ALIGN.LEFT)
    for run in paragraph.runs:
        font = run.font
        font.name = style.get("font_family", "Microsoft YaHei")
        font.size = Pt(float(style.get("font_size", max(8, element["height"] * 0.72))))
        font.bold = bool(style.get("bold", False))
        font.italic = bool(style.get("italic", False))
        rgb = hex_to_rgb(style.get("font_color", "#000000"))
        font.color.rgb = RGBColor(*rgb)


def apply_text_frame(shape: Any, element: dict[str, Any], style_key: str, helpers: tuple[Any, ...]) -> None:
    _, _, _, PP_ALIGN, MSO_ANCHOR, _, _, _, _ = helpers
    text = element.get("text", element.get("content", ""))
    runs = element.get("runs")
    if not text and not runs:
        return
    style = element.get(style_key, element.get("style", {}))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = px(float(style.get("margin_left", 0)))
    frame.margin_right = px(float(style.get("margin_right", 0)))
    frame.margin_top = px(float(style.get("margin_top", 0)))
    frame.margin_bottom = px(float(style.get("margin_bottom", 0)))
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(style.get("vertical_alignment", "top")).lower(), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(style.get("alignment", "left")).lower(), PP_ALIGN.LEFT)
    if runs:
        for run_data in runs:
            run_style = {**style, **run_data}
            run = paragraph.add_run()
            run.text = str(run_data.get("text", ""))
            style_run(run, run_style, helpers)
    else:
        run = paragraph.add_run()
        run.text = str(text)
        style_run(run, style, helpers)


def add_shape(slide: Any, element: dict[str, Any], helpers: tuple[Any, ...]) -> None:
    _, RGBColor, MSO_SHAPE, _, _, _, _, _, _ = helpers
    shape_name = str(element.get("shape", "rect")).lower()
    shape_type = {
        "ellipse": MSO_SHAPE.OVAL,
        "line": MSO_SHAPE.LINE_INVERSE,
        "rect": MSO_SHAPE.RECTANGLE,
    }.get(shape_name, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, px(element["x"]), px(element["y"]), px(element["width"]), px(element["height"]))
    style = element.get("style", {})
    if "fill_color" in style:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(style["fill_color"]))
    if "line_color" in style:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(style["line_color"]))
    if "line_width" in style:
        shape.line.width = px(float(style["line_width"]))
    apply_text_frame(shape, element, "text_style", helpers)


def apply_cell_border(cell: Any, color: str, width: float, helpers: tuple[Any, ...]) -> None:
    _, _, _, _, _, _, parse_xml, nsdecls, qn = helpers
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    width_emu = max(12700, px(width))
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tc_pr.find(qn(edge))
        if existing is not None:
            tc_pr.remove(existing)
        fragment = (
            f"<{edge} {nsdecls('a')} w=\"{width_emu}\">"
            f"<a:solidFill><a:srgbClr val=\"{color.strip().lstrip('#')}\"/></a:solidFill>"
            "<a:prstDash val=\"solid\"/>"
            f"</{edge}>"
        )
        tc_pr.append(parse_xml(fragment))


def style_run(run: Any, style: dict[str, Any], helpers: tuple[Any, ...]) -> None:
    _, RGBColor, _, _, _, Pt, _, _, _ = helpers
    font = run.font
    font.name = style.get("font_family", "Microsoft YaHei")
    font.size = Pt(float(style.get("font_size", 14)))
    font.bold = bool(style.get("bold", False))
    font.italic = bool(style.get("italic", False))
    font.color.rgb = RGBColor(*hex_to_rgb(style.get("font_color", "#000000")))


def apply_cell_text(cell: Any, cell_data: dict[str, Any], base_style: dict[str, Any], helpers: tuple[Any, ...]) -> None:
    _, _, _, PP_ALIGN, MSO_ANCHOR, _, _, _, _ = helpers
    style = {**base_style, **cell_data.get("style", {})}
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.margin_left = px(float(style.get("margin_left", 4)))
    text_frame.margin_right = px(float(style.get("margin_right", 4)))
    text_frame.margin_top = px(float(style.get("margin_top", 1)))
    text_frame.margin_bottom = px(float(style.get("margin_bottom", 1)))
    text_frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(style.get("vertical_alignment", "middle")).lower(), MSO_ANCHOR.MIDDLE)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(style.get("alignment", "left")).lower(), PP_ALIGN.LEFT)

    runs = cell_data.get("runs")
    if runs:
        for run_data in runs:
            run_style = {**style, **run_data}
            run = paragraph.add_run()
            run.text = str(run_data.get("text", ""))
            style_run(run, run_style, helpers)
    else:
        run = paragraph.add_run()
        run.text = str(cell_data.get("text", ""))
        style_run(run, style, helpers)


def add_table(slide: Any, element: dict[str, Any], helpers: tuple[Any, ...]) -> None:
    _, RGBColor, _, _, _, _, _, _, _ = helpers
    rows = element.get("rows", [])
    columns = element.get("columns", [])
    row_count = len(rows)
    col_count = len(columns)
    if row_count == 0 or col_count == 0:
        return

    shape = slide.shapes.add_table(
        row_count,
        col_count,
        px(element["x"]),
        px(element["y"]),
        px(element["width"]),
        px(element["height"]),
    )
    table = shape.table
    for index, column in enumerate(columns):
        table.columns[index].width = px(float(column.get("width", element["width"] / col_count)))
    for index, row in enumerate(rows):
        table.rows[index].height = px(float(row.get("height", element["height"] / row_count)))

    base_style = {
        "font_family": "Microsoft YaHei",
        "font_size": 14,
        "font_color": "#000000",
        "bold": False,
        "fill_color": "#FFFFFF",
        "line_color": "#000000",
        "line_width": 1,
        **element.get("style", {}),
    }

    for row_idx, row in enumerate(rows):
        for col_idx, cell_data in enumerate(row.get("cells", [])):
            if col_idx >= col_count:
                continue
            cell = table.cell(row_idx, col_idx)
            row_span = int(cell_data.get("row_span", 1))
            col_span = int(cell_data.get("col_span", 1))
            if row_span > 1 or col_span > 1:
                end_row = min(row_idx + row_span - 1, row_count - 1)
                end_col = min(col_idx + col_span - 1, col_count - 1)
                cell.merge(table.cell(end_row, end_col))
            style = {**base_style, **cell_data.get("style", {})}
            if "fill_color" in style:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*hex_to_rgb(style["fill_color"]))
            apply_cell_border(cell, style.get("line_color", "#000000"), float(style.get("line_width", 1)), helpers)
            apply_cell_text(cell, cell_data, base_style, helpers)


def build_pptx(scene_path: Path, out_path: Path) -> None:
    helpers = require_pptx()
    Presentation = helpers[0]
    scene = json.loads(scene_path.read_text(encoding="utf-8"))
    scene_dir = scene_path.parent

    prs = Presentation()
    prs.slide_width = px(scene["canvas"]["width"])
    prs.slide_height = px(scene["canvas"]["height"])
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = scene.get("background", {})
    if background.get("type") == "solid":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = helpers[1](*hex_to_rgb(background.get("color", "#FFFFFF")))
    elif background.get("type") == "image" and background.get("path"):
        bg_path = resolve_path(scene_dir, background["path"])
        slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    for element in sorted(scene.get("elements", []), key=lambda item: item.get("z_index", 0)):
        kind = element.get("type")
        if kind == "image":
            path = resolve_path(scene_dir, element["path"])
            slide.shapes.add_picture(
                str(path),
                px(element["x"]),
                px(element["y"]),
                width=px(element["width"]),
                height=px(element["height"]),
            )
        elif kind == "text":
            add_text(slide, element, helpers)
        elif kind == "shape":
            add_shape(slide, element, helpers)
        elif kind == "table":
            add_table(slide, element, helpers)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPTX from scene JSON.")
    parser.add_argument("--scene", required=True, help="Path to scene.json.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    args = parser.parse_args()
    build_pptx(Path(args.scene).resolve(), Path(args.out).resolve())


if __name__ == "__main__":
    main()
